from __future__ import annotations

import json
from io import BytesIO
from pathlib import Path

import pytest
from openpyxl import load_workbook
from pptx import Presentation

from opensquilla.artifacts import ArtifactStore
from opensquilla.tools.builtin.file_authoring import (
    create_csv,
    create_pdf_report,
    create_pptx,
    create_xlsx,
)
from opensquilla.tools.types import CallerKind, ToolContext, current_tool_context


def _channel_artifact_context(tmp_path: Path) -> ToolContext:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    return ToolContext(
        is_owner=False,
        caller_kind=CallerKind.CHANNEL,
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:feishu:group:oc_demo",
    )


def _published_material(ctx: ToolContext, result: str) -> tuple[dict[str, object], bytes]:
    payload = json.loads(result)
    assert payload["status"] == "published"
    assert "download_url" not in payload["artifact"]
    assert len(ctx.published_artifacts) == 1
    assert ctx.published_artifacts[0]["download_url"].startswith("/api/v1/artifacts/")

    store = ArtifactStore(ctx.artifact_media_root or "")
    _, path = store.resolve_for_download(
        str(payload["artifact"]["id"]),
        session_id=str(payload["artifact"]["session_id"]),
    )
    return payload["artifact"], path.read_bytes()


@pytest.mark.asyncio
async def test_create_csv_publishes_channel_artifact(tmp_path: Path) -> None:
    ctx = _channel_artifact_context(tmp_path)
    token = current_tool_context.set(ctx)
    try:
        result = await create_csv(
            name="usage.csv",
            rows=[["name", "count"], ["alpha", 2], ["beta", 3]],
        )
    finally:
        current_tool_context.reset(token)

    artifact, material = _published_material(ctx, result)
    assert artifact["name"] == "usage.csv"
    assert artifact["mime"] == "text/csv"
    assert material.decode("utf-8-sig").splitlines() == [
        "name,count",
        "alpha,2",
        "beta,3",
    ]


@pytest.mark.asyncio
async def test_create_xlsx_publishes_channel_artifact(tmp_path: Path) -> None:
    ctx = _channel_artifact_context(tmp_path)
    token = current_tool_context.set(ctx)
    try:
        result = await create_xlsx(
            name="metrics.xlsx",
            sheets=[
                {
                    "name": "Summary",
                    "rows": [["metric", "value"], ["requests", 42]],
                }
            ],
        )
    finally:
        current_tool_context.reset(token)

    artifact, material = _published_material(ctx, result)
    assert artifact["name"] == "metrics.xlsx"
    assert artifact["mime"] == (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
    assert material.startswith(b"PK")
    workbook = load_workbook(BytesIO(material))
    sheet = workbook["Summary"]
    assert sheet["A1"].value == "metric"
    assert sheet["B2"].value == 42


@pytest.mark.asyncio
async def test_create_pptx_publishes_channel_artifact(tmp_path: Path) -> None:
    ctx = _channel_artifact_context(tmp_path)
    token = current_tool_context.set(ctx)
    try:
        result = await create_pptx(
            name="brief.pptx",
            slides=[
                {
                    "title": "Launch Readiness",
                    "bullets": ["Group reply works", "File artifacts are safe"],
                }
            ],
        )
    finally:
        current_tool_context.reset(token)

    artifact, material = _published_material(ctx, result)
    assert artifact["name"] == "brief.pptx"
    assert artifact["mime"] == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert material.startswith(b"PK")
    presentation = Presentation(BytesIO(material))
    slide = presentation.slides[0]
    assert slide.shapes.title.text == "Launch Readiness"
    assert "Group reply works" in slide.placeholders[1].text


@pytest.mark.asyncio
async def test_create_pdf_report_publishes_channel_artifact(tmp_path: Path) -> None:
    ctx = _channel_artifact_context(tmp_path)
    token = current_tool_context.set(ctx)
    try:
        result = await create_pdf_report(
            name="report.pdf",
            title="Channel File Delivery & Safety",
            sections=[
                {
                    "heading": "Result <Verified>",
                    "body": "Generated files are published as channel artifacts: A & B.",
                }
            ],
        )
    finally:
        current_tool_context.reset(token)

    artifact, material = _published_material(ctx, result)
    assert artifact["name"] == "report.pdf"
    assert artifact["mime"] == "application/pdf"
    assert material.startswith(b"%PDF")
